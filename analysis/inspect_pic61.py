"""Extract Picture 61 from slide 7 and list ALL shapes including group children."""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
OUT_DIR = r'D:\Projects\Sales_Platform\analysis'

prs = Presentation(SRC)
slide7 = prs.slides[6]

print('=== SLIDE 7 ALL SHAPES (including groups) ===')

def print_shapes(shapes, indent=0):
    prefix = '  ' * indent
    for s in shapes:
        l = round(s.left/914400, 2) if s.left else '?'
        t = round(s.top/914400, 2) if s.top else '?'
        w = round(s.width/914400, 2) if s.width else '?'
        h = round(s.height/914400, 2) if s.height else '?'
        txt = ''
        if hasattr(s, 'text') and s.text:
            txt = f'  text="{s.text[:30]}"'
        print(f'{prefix}type={s.shape_type}({int(s.shape_type)})  name="{s.name}"  x={l}  y={t}  w={w}  h={h}{txt}')
        # Recurse into groups
        if int(s.shape_type) == 6:  # GROUP_SHAPE
            print_shapes(s.shapes, indent+1)

print_shapes(slide7.shapes)

print()
print('=== PICTURES ON SLIDE 7 ===')
pic_count = 0
for s in slide7.shapes:
    if int(s.shape_type) == 13:  # PICTURE
        pic_count += 1
        l = round(s.left/914400, 2)
        t = round(s.top/914400, 2)
        w = round(s.width/914400, 2)
        print(f'  name="{s.name}"  x={l}  y={t}  w={w}')
        # Extract image
        try:
            img_blob = s.image.blob
            ext = s.image.ext
            out_path = os.path.join(OUT_DIR, f'pic_{s.name.replace(" ","_")}_slide7.{ext}')
            with open(out_path, 'wb') as f:
                f.write(img_blob)
            print(f'  → Saved: {out_path} ({len(img_blob)//1024}KB)')
        except Exception as e:
            print(f'  → Image extract error: {e}')

print(f'Total pictures: {pic_count}')

print()
print('=== CHARTS ON SLIDE 7 (via rels) ===')
for rel in slide7.part.rels.values():
    if 'chart' in rel.reltype.lower():
        print(f'  rel: {rel.reltype}')
        try:
            tgt = rel._target
            if hasattr(tgt, 'blob'):
                blob = tgt.blob
                tree = etree.fromstring(blob)
                # Get series
                C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
                for ser in tree.findall(f'.//{{{C}}}ser'):
                    name_el = ser.find(f'.//{{{C}}}v')
                    cats = [c.get('v','') for c in ser.findall(f'.//{{{C}}}cat//{{{C}}}v')]
                    sname = name_el.text if name_el is not None else '?'
                    print(f'    series="{sname}" categories={cats[:5]}')
        except Exception as e:
            print(f'  → Error: {e}')
