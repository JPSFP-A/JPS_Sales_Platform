"""
Fix slide numbers: remove non-rendering master field, add explicit numbers to every slide.
Dark slides (cover, section dividers, thank you) → white text
Content slides → Oxford Blue text
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
OUT = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'

P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

OX = RGBColor(0x06, 0x25, 0x52)
WH = RGBColor(0xFF, 0xFF, 0xFF)

# 1-indexed dark slides (section dividers): white text
DARK = {4, 11, 15, 17, 20}
# Skip cover and thank you — no page number convention
SKIP = {1, 34}

# x/y/w/h for slide number box — bottom right, aligned with footer
SN_X = Inches(11.5)
SN_Y = Inches(7.12)
SN_W = Inches(1.6)
SN_H = Inches(0.28)

prs = Presentation(SRC)

# ── 1. Remove MasterSlideNum shape from master (if present) ──────────────────
master = prs.slide_masters[0]
cSld = master._element.find(f'{{{P}}}cSld')
spTree = cSld.find(f'{{{P}}}spTree')
removed = 0
for sp in list(spTree):
    cNvPr = sp.find(f'.//{{{P}}}cNvPr')
    if cNvPr is None:
        cNvPr = sp.find(f'.//{{{A}}}cNvPr')
    name = (cNvPr.get('name', '') if cNvPr is not None else '')
    if name == 'MasterSlideNum':
        spTree.remove(sp)
        removed += 1
        continue
    # Also remove old ph-based sldNum placeholders if any remain
    ph = sp.find(f'.//{{{P}}}ph')
    if ph is not None and ph.get('type') == 'sldNum':
        spTree.remove(sp)
        removed += 1
print(f'Removed {removed} master sldNum shape(s) ✓')

# ── 2. Remove stale explicit slide-number boxes from slides (re-run safety) ──
SLIDENUM_TAG = 'SlideNumExplicit'
for slide in prs.slides:
    for shape in list(slide.shapes):
        if shape.name == SLIDENUM_TAG:
            sp_elem = shape._element
            sp_elem.getparent().remove(sp_elem)

# ── 3. Add explicit slide number text box to every slide ─────────────────────
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in SKIP:
        print(f'  Slide {num:2d} (skip)')
        continue
    dark = num in DARK
    color = WH if dark else OX

    txBox = slide.shapes.add_textbox(SN_X, SN_Y, SN_W, SN_H)
    txBox.name = SLIDENUM_TAG

    tf = txBox.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = str(num)
    run.font.name = 'Aptos'
    run.font.size = Pt(11)
    run.font.bold = False
    run.font.color.rgb = color

    label = 'dark' if dark else 'white'
    print(f'  Slide {num:2d} ({label}) → "{num}"')

prs.save(OUT)
print(f'\nSaved → {OUT}')
print(f'Size: {round(os.path.getsize(OUT)/1024)}KB')
