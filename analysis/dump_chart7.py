"""Dump full chart XML for slide 7 to understand legend position."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

prs = Presentation(SRC)
slide7 = prs.slides[6]

for shape in slide7.shapes:
    if shape.shape_type != 3:
        continue
    chart_tree = etree.fromstring(shape.chart.part.blob)
    full = etree.tostring(chart_tree, pretty_print=True).decode()
    # Show legend section
    print('=== LEGEND ===')
    for leg in chart_tree.findall(f'.//{{{C}}}legend'):
        print(etree.tostring(leg, pretty_print=True).decode())
    print()
    # Show plot area
    print('=== PLOT AREA ===')
    for pa in chart_tree.findall(f'.//{{{C}}}plotArea'):
        print(etree.tostring(pa, pretty_print=True).decode()[:800])
    print()
    # Show outer chart layout
    print('=== CHART LAYOUT ===')
    for cl in chart_tree.findall(f'.//{{{C}}}chart/{{{C}}}plotVisOnly'):
        print(etree.tostring(cl, pretty_print=True).decode())
    print()
    print('=== ALL DIRECT CHART CHILDREN ===')
    chart_el = chart_tree.find(f'{{{C}}}chart')
    if chart_el is not None:
        for child in chart_el:
            print(f'  {child.tag}', child.attrib)
    print()
    print('=== FULL XML (first 5000 chars) ===')
    print(full[:5000])
