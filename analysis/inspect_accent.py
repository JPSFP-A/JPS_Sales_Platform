"""
Find thin horizontal accent shapes near title on each content slide.
Also dump chart XML for slide 7.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
prs = Presentation(SRC)

DARK = {1, 4, 11, 15, 17, 20, 34}

print('=== ACCENT LINES (thin horizontal shapes y<0.9") ===')
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in DARK or num in {1, 34}:
        continue
    for s in slide.shapes:
        if s.left is None or s.top is None:
            continue
        t = s.top / 914400
        w = s.width / 914400
        h = s.height / 914400
        x = s.left / 914400
        if t < 0.9 and h < 0.12 and w > 0.1:
            # get fill color
            clr = '?'
            try:
                fc = s.fill.fore_color
                clr = str(fc.rgb) if fc.type else 'none'
            except:
                pass
            print(f'  Slide {num:2d}  name="{s.name}"  x={round(x,2)}"  y={round(t,2)}"  '
                  f'w={round(w,2)}"  h={round(h,3)}"  fill={clr}')

print()
print('=== SLIDE 7 CHART XML (legend section) ===')
slide7 = prs.slides[6]
for s in slide7.shapes:
    if s.shape_type == 3:  # CHART
        chart_part = s.part
        chart_xml = chart_part.blob
        tree = etree.fromstring(chart_xml)
        # Find legend element
        nsmap = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        legends = tree.findall('.//c:legend', nsmap)
        for leg in legends:
            print(etree.tostring(leg, pretty_print=True).decode())
        # Also find plotArea layout
        plots = tree.findall('.//c:plotArea/c:layout', nsmap)
        for p in plots:
            print('plotArea layout:', etree.tostring(p, pretty_print=True).decode())
        print('(full chart XML first 3000 chars):')
        print(etree.tostring(tree, pretty_print=True).decode()[:3000])
