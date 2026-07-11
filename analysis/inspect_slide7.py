"""Inspect all shapes on slide 7 and dump title XML for slides with odd coloring."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
prs = Presentation(SRC)

slide7 = prs.slides[6]  # 0-indexed
print('=== SLIDE 7 ALL SHAPES ===')
for s in slide7.shapes:
    l = round(s.left/914400, 2) if s.left else '?'
    t = round(s.top/914400, 2) if s.top else '?'
    w = round(s.width/914400, 2) if s.width else '?'
    h = round(s.height/914400, 2) if s.height else '?'
    txt = ''
    if hasattr(s, 'text') and s.text:
        txt = f'  text="{s.text[:40]}"'
    print(f'  type={s.shape_type}  name="{s.name}"  x={l}  y={t}  w={w}  h={h}{txt}')

# Dump title XML for slides where title has old colors
print()
print('=== TITLE XML for slides with per-word color (7, 8, 21) ===')
for idx in [6, 7, 20]:  # 0-indexed
    slide = prs.slides[idx]
    print(f'\n--- Slide {idx+1} title shapes ---')
    for s in slide.shapes:
        if s.top is not None and s.top < Inches(1.5) and hasattr(s, 'text_frame'):
            txt = s.text[:60] if s.text else ''
            if txt and len(txt) > 3:
                print(f'  Shape: {s.name}  text="{txt}"')
                # Dump the XML
                xml = etree.tostring(s._element, pretty_print=True).decode()
                print(xml[:2000])
                print('...')
                break
