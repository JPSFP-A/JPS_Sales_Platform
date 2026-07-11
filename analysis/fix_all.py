"""
Three-pass update to May_2026_Corp_Perf_JPS_Brand.pptx:
  1. Slide master: white background + JPS theme colors + Aptos fonts
  2. Slide numbers: add sldNum field to master spTree (inherits to all slides)
  3. TOC hyperlinks: verify/fix slide targets for all 6 entries
"""
import sys, uuid, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
OUT = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'

A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
SL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'

def emu(inches):
    return str(int(inches * 914400))

prs = Presentation(SRC)
master = prs.slide_masters[0]

# ── 1. MASTER BACKGROUND → white solid ───────────────────────────────────────
cSld = master._element.find(f'{{{P}}}cSld')
bg   = cSld.find(f'{{{P}}}bg')
bgPr = bg.find(f'{{{P}}}bgPr') if bg is not None else None

if bgPr is not None:
    for child in list(bgPr):
        bgPr.remove(child)
    solid = etree.SubElement(bgPr, f'{{{A}}}solidFill')
    srgb  = etree.SubElement(solid, f'{{{A}}}srgbClr')
    srgb.set('val', 'FFFFFF')
    etree.SubElement(bgPr, f'{{{A}}}effectLst')
    print('Master background → white solid ✓')
else:
    # Build bg from scratch
    if bg is None:
        bg = etree.SubElement(cSld, f'{{{P}}}bg')
    bgPr = etree.SubElement(bg, f'{{{P}}}bgPr')
    solid = etree.SubElement(bgPr, f'{{{A}}}solidFill')
    srgb  = etree.SubElement(solid, f'{{{A}}}srgbClr')
    srgb.set('val', 'FFFFFF')
    etree.SubElement(bgPr, f'{{{A}}}effectLst')
    print('Master background → white solid (created) ✓')

# ── 2. THEME: JPS colors + Aptos fonts ───────────────────────────────────────
theme_part = master.part.part_related_by(
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
tree = etree.fromstring(theme_part.blob)

clr_scheme = tree.find(f'.//{{{A}}}clrScheme')
COLOR_MAP = {
    'dk1':      ('srgbClr', '062552'),
    'dk2':      ('srgbClr', '0E4E95'),
    'lt2':      ('srgbClr', 'F2F2F2'),
    'accent1':  ('srgbClr', '009FDA'),
    'accent2':  ('srgbClr', 'FFF000'),
    'accent3':  ('srgbClr', '062552'),
    'accent4':  ('srgbClr', '0E4E95'),
    'accent5':  ('srgbClr', '009FDA'),
    'accent6':  ('srgbClr', '44546A'),
    'hlink':    ('srgbClr', '0E4E95'),
    'folHlink': ('srgbClr', '062552'),
}
for tag, (kind, val) in COLOR_MAP.items():
    el = clr_scheme.find(f'{{{A}}}{tag}')
    if el is None:
        continue
    for child in list(el):
        el.remove(child)
    sc = etree.SubElement(el, f'{{{A}}}{kind}')
    sc.set('val', val)
clr_scheme.set('name', 'JPS 2026')
print('Theme colors → JPS 2026 ✓')

font_scheme = tree.find(f'.//{{{A}}}fontScheme')
major = font_scheme.find(f'{{{A}}}majorFont')
minor = font_scheme.find(f'{{{A}}}minorFont')
major.find(f'{{{A}}}latin').set('typeface', 'Aptos Black')
minor.find(f'{{{A}}}latin').set('typeface', 'Aptos')
font_scheme.set('name', 'JPS 2026')
print('Theme fonts → Aptos Black / Aptos ✓')

theme_part._blob = etree.tostring(tree)

# ── 3. SLIDE NUMBER → regular master shape (no ph element) ───────────────────
# Regular shapes in the master spTree appear on ALL slides unconditionally.
# Placeholder shapes (ph) only resolve via layout inheritance — won't render
# on slides that don't have a matching layout placeholder.
# Position: x=11.5" y=7.1" w=1.6" h=0.28" (bottom-right, 11pt Oxford Blue)
SN_XML = f'''<p:sp xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">
  <p:nvSpPr>
    <p:cNvPr id="901" name="MasterSlideNum"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{emu(11.5)}" y="{emu(7.12)}"/>
      <a:ext cx="{emu(1.63)}" cy="{emu(0.28)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" rtlCol="0"/>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="r"/>
      <a:fld id="{{{str(uuid.uuid4()).upper()}}}" type="slidenum">
        <a:rPr lang="en-US" sz="1000" b="0" smtClean="0">
          <a:solidFill><a:srgbClr val="062552"/></a:solidFill>
          <a:latin typeface="Aptos"/>
        </a:rPr>
        <a:t>‹#›</a:t>
      </a:fld>
    </a:p>
  </p:txBody>
</p:sp>'''

spTree = cSld.find(f'{{{P}}}spTree')
# Remove any previous MasterSlideNum shapes
for sp in list(spTree):
    cNvPr = sp.find(f'.//{{{P}}}cNvPr') or sp.find(f'.//{{{A}}}cNvPr')
    if cNvPr is None:
        cNvPr = sp.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    if cNvPr is not None and cNvPr.get('name', '') == 'MasterSlideNum':
        spTree.remove(sp)
    # Also remove old ph-based sldNum
    ph = sp.find(f'.//{{{P}}}ph')
    if ph is not None and ph.get('type') == 'sldNum':
        spTree.remove(sp)
spTree.append(etree.fromstring(SN_XML))
print('Slide number (regular master shape) added ✓')

# ── 4. FIX TOC HYPERLINKS ────────────────────────────────────────────────────
TOC_MAP = {
    0: 3,   # Executive Summary
    1: 4,   # Financial Highlights
    2: 11,  # Operational Expenditure
    3: 15,  # Capital Expenditure
    4: 17,  # Cash Flow and Funding Outlook
    5: 20,  # Appendices
}

slide2      = prs.slides[1]
slide2_part = slide2.part

# Ensure relationships exist for all target slides
target_rids = {}
for para_idx, slide_num in TOC_MAP.items():
    target_part = prs.slides[slide_num - 1].part
    # Find existing rel
    rid = None
    for r_id, rel in slide2_part.rels.items():
        try:
            if rel._target.partname == target_part.partname:
                rid = r_id
                break
        except AttributeError:
            pass
    if not rid:
        rid = slide2_part.relate_to(target_part, SL)
    target_rids[slide_num] = rid
    print(f'  Slide {slide_num:2d} → {rid}')

# Apply hlinkClick to correct page-number runs
for shape in slide2.shapes:
    if 'Content' in shape.name and shape.has_text_frame:
        for i, para in enumerate(shape.text_frame.paragraphs):
            if i not in TOC_MAP:
                continue
            target = TOC_MAP[i]
            rid    = target_rids[target]
            for run in para.runs:
                # Clean stray hlinkClick on any non-target run
                rpr = run._r.find(f'{{{A}}}rPr')
                if rpr is not None:
                    existing = rpr.find(f'{{{A}}}hlinkClick')
                    if existing is not None and run.text.strip() != str(target):
                        rpr.remove(existing)
                # Add to the exact page-number run
                if run.text.strip() == str(target):
                    if rpr is None:
                        rpr = etree.SubElement(run._r, f'{{{A}}}rPr')
                        run._r.insert(0, rpr)
                    # Remove old hlinkClick
                    old = rpr.find(f'{{{A}}}hlinkClick')
                    if old is not None:
                        rpr.remove(old)
                    # Add fresh hlinkClick
                    hc = etree.SubElement(rpr, f'{{{A}}}hlinkClick')
                    hc.set(f'{{{R}}}id', rid)
                    hc.set('action', 'ppaction://hlinksldjump')
                    print(f'  Para {i}: "{run.text}" → slide {target} ({rid}) ✓')

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'\nSaved → {OUT}')
print(f'Size: {round(os.path.getsize(OUT)/1024)}KB')
