import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

prs = Presentation(r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx')
slide = prs.slides[1]  # slide 2 (0-indexed)

# Find the TOC content placeholder
for shape in slide.shapes:
    if 'Content Placeholder' in shape.name or 'Placeholder' in shape.name:
        print(f'Shape: {shape.name}')
        xml = etree.tostring(shape._element, pretty_print=True).decode()
        print(xml[:4000])
        print('...')
