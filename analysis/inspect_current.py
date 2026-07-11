import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

prs = Presentation(r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx')
print(f'Slides: {len(prs.slides)}')

# Check slide 2 (likely TOC) in detail
for idx in [1, 2]:
    slide = prs.slides[idx]
    print(f'\n=== Slide {idx+1} ===')
    for s in slide.shapes:
        t = ''
        if s.has_text_frame:
            t = s.text_frame.text[:200].replace('\n',' | ').strip()
        print(f'  {s.name} [{s.shape_type}]: {t}')

# Check master background and theme
master = prs.slide_masters[0]
bg = master.background
print(f'\nMaster bg fill type: {bg.fill.type}')
print(f'Master layouts: {[l.name for l in master.slide_layouts]}')

# Check if slide numbers enabled in prs XML
prs_xml_str = etree.tostring(prs.element, pretty_print=True).decode()
# Look for showPr or slideNumPlaceholders
import re
sn_refs = re.findall(r'<[^>]*sldNum[^>]*>', prs_xml_str)
print(f'\nSlide number refs in prs XML: {sn_refs[:5]}')

# Check if any layout has sldNum placeholder
for layout in master.slide_layouts:
    for sp in layout.shapes:
        ph = sp.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}ph')
        if ph is not None and ph.get('type') == 'sldNum':
            print(f'  Layout {layout.name!r} has sldNum placeholder: {sp.name}')
