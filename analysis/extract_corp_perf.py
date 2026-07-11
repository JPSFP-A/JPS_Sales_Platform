import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'C:\Users\jwilson\Downloads\May 2026 Corporate Performance Report.pptx')
print(f'Slide size: {prs.slide_width.inches:.3f}" x {prs.slide_height.inches:.3f}"')
print(f'Slide count: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides, 1):
    print(f'=== SLIDE {i} ===')
    for shape in slide.shapes:
        stype = shape.shape_type
        nm = shape.name
        if shape.has_text_frame:
            txt = ' / '.join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
            if txt:
                print(f'  TEXT [{nm}]: {txt[:300]}')
        if shape.has_table:
            tbl = shape.table
            rows_data = []
            for r in tbl.rows:
                row_cells = [c.text_frame.text.strip().replace('\n',' ') for c in r.cells]
                rows_data.append(' | '.join(row_cells))
            print(f'  TABLE [{nm}] ({len(tbl.rows)}r x {len(tbl.columns)}c):')
            for r in rows_data[:6]:
                print(f'    {r[:180]}')
            if len(rows_data) > 6:
                print(f'    ... ({len(rows_data)-6} more rows)')
        if stype == MSO_SHAPE_TYPE.CHART:
            ch = shape.chart
            print(f'  CHART [{nm}]: type={ch.chart_type}')
            try:
                for series in ch.series:
                    vals = list(series.values)[:8]
                    print(f'    series "{series.name}": {vals}')
            except Exception as e:
                print(f'    (series read error: {e})')
        if stype == MSO_SHAPE_TYPE.PICTURE:
            print(f'  IMAGE [{nm}]')
    print()
