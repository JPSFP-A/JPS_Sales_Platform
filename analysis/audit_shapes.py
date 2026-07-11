"""
Audit all content slides:
- Fonts used (flag non-Aptos/non-brand fonts)
- Title shape text + underline attributes per run
- Shapes in logo zone (x > 10.5", y < 1.5")
- All shape types per slide
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'

# Dark slides we skip in content processing
DARK = {1, 4, 11, 15, 17, 20, 34}
SKIP = {1, 34}

# Logo zone: x > 10.5", y < 1.5" (anything that isn't our logo)
LOGO_X = Inches(10.5)
LOGO_Y = Inches(1.5)

prs = Presentation(SRC)

print('=' * 70)
print('FONT AUDIT — non-brand fonts in text shapes')
print('=' * 70)
font_issues = {}
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in SKIP or num in DARK:
        continue
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                fname = run.font.name
                if fname and fname.lower() not in ('aptos', 'aptos black', 'calibri', 'none', ''):
                    key = (num, shape.name, fname)
                    txt = run.text[:40].replace('\n', ' ')
                    font_issues[key] = font_issues.get(key, [])
                    font_issues[key].append(txt)

for (snum, sname, fname), samples in sorted(font_issues.items()):
    print(f'  Slide {snum:2d}  [{fname}]  shape="{sname}"')
    for s in samples[:3]:
        print(f'           → "{s}"')

print()
print('=' * 70)
print('TITLE UNDERLINE AUDIT — per-run underline/color')
print('=' * 70)
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in SKIP or num in DARK:
        continue
    for shape in slide.shapes:
        # Title shapes are typically at top of slide or named "Title"
        if shape.shape_type not in (1, 13, 14, 17, 19):  # text-bearing types
            continue
        if not shape.has_text_frame:
            continue
        if shape.top > Inches(1.8):  # skip shapes below title area
            continue
        tf = shape.text_frame
        full_text = tf.text.strip()
        if not full_text or len(full_text) < 3:
            continue
        for pi, para in enumerate(tf.paragraphs):
            for ri, run in enumerate(para.runs):
                u = run.font.underline
                clr = None
                try:
                    clr = run.font.color.rgb
                except:
                    pass
                if u or (clr and str(clr) not in ('062552', '000000', 'FFFFFF', 'None')):
                    print(f'  Slide {num:2d}  shape="{shape.name}"  run[{pi},{ri}]="{run.text}"  '
                          f'underline={u}  color={clr}')

print()
print('=' * 70)
print('LOGO ZONE AUDIT — shapes in top-right (x>10.5", y<1.5")')
print('=' * 70)
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in SKIP or num in DARK:
        continue
    for shape in slide.shapes:
        try:
            lx = shape.left
            ty = shape.top
        except:
            continue
        if lx is not None and ty is not None:
            if lx > LOGO_X and ty < LOGO_Y:
                print(f'  Slide {num:2d}  type={shape.shape_type}  name="{shape.name}"  '
                      f'left={round(lx/914400,2)}"  top={round(ty/914400,2)}"  '
                      f'w={round(shape.width/914400,2)}"')
