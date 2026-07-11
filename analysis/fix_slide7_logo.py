"""Replace slide 7 Picture 61 (3KB chart image) with the JPS logo from slide 1."""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'

prs = Presentation(SRC)

# ── Get JPS logo blob from slide 1 ─────────────────────────────────────────
logo_blob = None
logo_content_type = None
slide1 = prs.slides[0]
for s in slide1.shapes:
    if int(s.shape_type) != 13:
        continue
    lx = s.left / 914400
    ty = s.top / 914400
    if lx > 10.0 and ty < 1.5:
        logo_blob = s.image.blob
        logo_content_type = s.image.content_type
        print(f'Logo from slide 1: "{s.name}"  {len(logo_blob)//1024}KB  type={logo_content_type}')
        break

if not logo_blob:
    print('ERROR: logo not found on slide 1')
    sys.exit(1)

# ── Replace Picture 61 blob on slide 7 ─────────────────────────────────────
slide7 = prs.slides[6]
fixed = False
for s in slide7.shapes:
    if s.name != 'Picture 61':
        continue
    pic_elem = s._pic
    blip = pic_elem.find('.//' + qn('a:blip'))
    if blip is None:
        print('ERROR: no blip element in Picture 61')
        break
    r_embed = blip.get(qn('r:embed'))
    img_part = s.part.related_part(r_embed)
    old_size = len(img_part.blob) if img_part.blob else 0
    img_part._blob = logo_blob
    # Also fix content type if needed
    if logo_content_type and img_part.content_type != logo_content_type:
        img_part._content_type = logo_content_type
    print(f'Replaced Picture 61: {old_size//1024}KB chart → {len(logo_blob)//1024}KB JPS logo')
    fixed = True
    break

if not fixed:
    print('ERROR: Picture 61 not found on slide 7')
    sys.exit(1)

prs.save(SRC)
print(f'Saved → {SRC}')
print(f'Size: {round(os.path.getsize(SRC)/1024)}KB')
