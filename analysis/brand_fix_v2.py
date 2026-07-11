"""
Comprehensive brand fix:
 1. Accent lines: 01ADEF → 009FDA (new Electric Blue) on all content slides
 2. Title text: color 002060 → 062552 (Oxford Blue), font → Century Gothic
 3. Body text: Arial → Montserrat (all text shapes)
 4. Slide 7 chart legend: remove manual layout (chart legend dragged to slide top-right)
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
OUT = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'

A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
C  = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
R  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

OX   = RGBColor(0x06, 0x25, 0x52)   # Oxford Blue (new brand)
EL   = RGBColor(0x00, 0x9F, 0xDA)   # Electric Blue (new brand)

DARK = {1, 4, 11, 15, 17, 20, 34}
SKIP = {1, 34}

# Shapes our scripts added — don't touch their fonts
OUR_SHAPES = {'SlideNumExplicit', 'MasterSlideNum'}

def is_our_shape(shape):
    return shape.name in OUR_SHAPES or shape.name.startswith('TextBox 6')

def fix_font(run, target_name):
    """Set run font name without touching color, size, bold, or underline."""
    rPr = run._r.find(f'{{{A}}}rPr')
    if rPr is None:
        return
    latin = rPr.find(f'{{{A}}}latin')
    if latin is None:
        latin = etree.SubElement(rPr, f'{{{A}}}latin')
    latin.set('typeface', target_name)

def fix_title_color(run):
    """Set run color to Oxford Blue in XML."""
    rPr = run._r.find(f'{{{A}}}rPr')
    if rPr is None:
        return
    # Remove any existing solidFill
    for sf in list(rPr.findall(f'{{{A}}}solidFill')):
        rPr.remove(sf)
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    sc = etree.SubElement(sf, f'{{{A}}}srgbClr')
    sc.set('val', '062552')

prs = Presentation(SRC)

# ── 1. ACCENT LINES: 01ADEF → 009FDA ────────────────────────────────────────
accent_fixed = 0
for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in SKIP or num in DARK:
        continue
    for s in slide.shapes:
        if s.left is None or s.top is None:
            continue
        t = s.top / 914400
        h = s.height / 914400
        if t < 0.9 and h < 0.12:
            try:
                clr = str(s.fill.fore_color.rgb)
                if clr in ('01ADEF', '0070C0', '009FDA', '003366'):
                    # Update via XML
                    spPr = s._element.find(f'{{{P}}}spPr') or s._element.find(f'{{{A}}}spPr')
                    if spPr is None:
                        spPr = s._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    # Find the srgbClr element and update
                    for srgb in s._element.iter(f'{{{A}}}srgbClr'):
                        parent = srgb.getparent()
                        grandparent = parent.getparent() if parent is not None else None
                        # Only update if this is a fill color (inside solidFill)
                        if parent is not None and 'solidFill' in parent.tag:
                            old = srgb.get('val', '')
                            if old.upper() in ('01ADEF', '0070C0', '003366'):
                                srgb.set('val', '009FDA')
                                accent_fixed += 1
            except Exception as e:
                pass

print(f'Accent lines updated: {accent_fixed} ✓')

# ── 2 & 3. TITLE + BODY FONT/COLOR ──────────────────────────────────────────
title_runs = 0
body_runs  = 0

for i, slide in enumerate(prs.slides):
    num = i + 1
    if num in SKIP or num in DARK:
        continue

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if is_our_shape(shape):
            continue
        if shape.left is None or shape.top is None:
            continue

        top_in = shape.top / 914400
        is_title = (top_in < 1.0 and shape.text_frame.text.strip())

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                old_font = run.font.name or ''
                is_arial = 'arial' in old_font.lower() if old_font else True

                if is_title:
                    # Fix font → Century Gothic
                    fix_font(run, 'Century Gothic')
                    # Fix color → Oxford Blue
                    try:
                        old_clr = str(run.font.color.rgb) if run.font.color.type else None
                        if old_clr in (None, '002060', '0A2F57', '002851', '003366', '1F3864'):
                            fix_title_color(run)
                    except:
                        fix_title_color(run)
                    title_runs += 1
                else:
                    # Body: only replace Arial with Montserrat
                    if is_arial or not old_font:
                        fix_font(run, 'Montserrat')
                        body_runs += 1

print(f'Title runs fixed (Century Gothic + Oxford Blue): {title_runs} ✓')
print(f'Body runs fixed (Montserrat): {body_runs} ✓')

# ── 4. SLIDE 7 CHART LEGEND: remove manual layout ───────────────────────────
slide7 = prs.slides[6]
legend_fixed = 0
for shape in slide7.shapes:
    if shape.shape_type != 3:  # CHART
        continue
    # Access chart XML via the graphic's relationship
    try:
        # Find the chart relationship id
        graphic_data = shape._element.find(
            './/{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData')
        if graphic_data is None:
            continue
        chart_elem = graphic_data.find(
            f'{{{C}}}chart')
        if chart_elem is None:
            continue
        r_id = chart_elem.get(f'{{{R}}}id')
        if not r_id:
            continue
        chart_part = slide7.part.related_parts.get(r_id)
        if chart_part is None:
            continue
        chart_tree = etree.fromstring(chart_part.blob)
        # Find and remove all manualLayout inside legend
        for ml in chart_tree.findall(f'.//{{{C}}}legend/{{{C}}}layout/{{{C}}}manualLayout'):
            parent = ml.getparent()
            if parent is not None:
                parent.remove(ml)
                legend_fixed += 1
        # Also check for legend with direct manualLayout
        for leg in chart_tree.findall(f'.//{{{C}}}legend'):
            for layout in leg.findall(f'{{{C}}}layout'):
                for ml in layout.findall(f'{{{C}}}manualLayout'):
                    layout.remove(ml)
                    legend_fixed += 1
        # Write back
        chart_part._blob = etree.tostring(chart_tree)
        print(f'Slide 7 chart legend manual layouts removed: {legend_fixed}')
    except Exception as e:
        print(f'Chart legend fix error: {e}')

if legend_fixed == 0:
    print('Slide 7: no manualLayout found in chart legend (legend may use default position)')

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'\nSaved → {OUT}')
print(f'Size: {round(os.path.getsize(OUT)/1024)}KB')
