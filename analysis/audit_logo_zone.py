"""Audit logo zone (x>10.5, y<1.5) across all slides — find misplaced images."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
prs = Presentation(SRC)

LOGO_X = 11.0   # shapes this far right
LOGO_Y = 1.5    # and this high up

print('=== LOGO ZONE AUDIT (x > 10.5", y < 1.5") ===')
print(f'{"Slide":>5}  {"Type":>6}  {"Name":<25}  {"x":>6}  {"y":>6}  {"w":>6}  {"Note"}')
print('-'*90)

for i, slide in enumerate(prs.slides):
    num = i + 1
    pics = []
    for s in slide.shapes:
        if s.left is None or s.top is None:
            continue
        lx = s.left / 914400
        ty = s.top / 914400
        if lx > LOGO_X and ty < LOGO_Y:
            tn = int(s.shape_type)
            w = round(s.width/914400, 2) if s.width else '?'
            note = ''
            if tn == 13:  # PICTURE
                try:
                    sz = len(s.image.blob)
                    note = f'image {sz//1024}KB'
                except:
                    note = 'image ?KB'
            elif tn == 17:  # TEXT_BOX
                note = f'text="{s.text[:20]}"'
            print(f'{num:>5}  {tn:>6}  {s.name:<25}  {round(lx,2):>6}  {round(ty,2):>6}  {w:>6}  {note}')
