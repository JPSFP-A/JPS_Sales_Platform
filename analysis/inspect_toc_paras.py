import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

prs = Presentation(r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx')
slide = prs.slides[1]  # slide 2

for shape in slide.shapes:
    if 'Content' in shape.name and shape.has_text_frame:
        print(f'Shape: {shape.name}')
        tf = shape.text_frame
        for i, para in enumerate(tf.paragraphs):
            print(f'\n  Para {i}: text={para.text[:120]!r}')
            for j, run in enumerate(para.runs):
                rpr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                has_link = False
                if rpr is not None:
                    hlink = rpr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick')
                    has_link = hlink is not None
                print(f'    Run {j}: {run.text!r}  has_link={has_link}')
