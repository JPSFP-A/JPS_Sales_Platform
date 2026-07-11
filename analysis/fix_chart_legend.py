"""Fix chart legend on slide 7 — legend was dragged to top-right of slide."""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

SRC = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
OUT = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'

C  = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

prs = Presentation(SRC)
slide7 = prs.slides[6]

fixed = 0
for shape in slide7.shapes:
    if shape.shape_type != 3:  # CHART
        continue
    print(f'Found chart: "{shape.name}" at x={round(shape.left/914400,2)}", y={round(shape.top/914400,2)}"')
    try:
        # Access chart via shape.chart (python-pptx API)
        chart_obj = shape.chart
        chart_part = chart_obj.part
        chart_tree = etree.fromstring(chart_part.blob)

        # Show current legend XML
        for leg in chart_tree.findall(f'.//{{{C}}}legend'):
            print('Legend XML:')
            print(etree.tostring(leg, pretty_print=True).decode())

            # Remove all manualLayout elements to reset to default position
            for ml in leg.findall(f'.//{{{C}}}manualLayout'):
                parent = ml.getparent()
                parent.remove(ml)
                fixed += 1
                print('Removed manualLayout ✓')

            # Also remove any layout wrapper that's now empty
            for lay in leg.findall(f'{{{C}}}layout'):
                if len(lay) == 0:
                    leg.remove(lay)

        if fixed == 0:
            print('No manualLayout in legend — checking all chart elements:')
            for ml in chart_tree.findall(f'.//{{{C}}}manualLayout'):
                parent = ml.getparent()
                gp = parent.getparent()
                print(f'  manualLayout inside {parent.tag} inside {gp.tag}')
                x_el = ml.find(f'{{{C}}}x')
                y_el = ml.find(f'{{{C}}}y')
                if x_el is not None:
                    print(f'    x={x_el.get("val")}  y={y_el.get("val") if y_el is not None else "?"}')

        # Write back
        chart_part._blob = etree.tostring(chart_tree)
        print(f'Chart part updated ({fixed} layouts removed)')
    except AttributeError as e:
        print(f'AttributeError: {e}')
        # Fallback: access via slide rels
        for rel in slide7.part.rels.values():
            try:
                tgt = rel._target
                if hasattr(tgt, 'blob'):
                    blob = tgt.blob
                    if blob and b'<c:chartSpace' in blob:
                        print(f'  Found chart via rels: {rel.reltype}')
                        tree = etree.fromstring(blob)
                        for ml in tree.findall(f'.//{{{C}}}manualLayout'):
                            parent = ml.getparent()
                            parent.remove(ml)
                            fixed += 1
                        tgt._blob = etree.tostring(tree)
            except:
                pass

prs.save(OUT)
print(f'\nSaved → {OUT}  ({round(os.path.getsize(OUT)/1024)}KB)')
print(f'Total fixed: {fixed}')
