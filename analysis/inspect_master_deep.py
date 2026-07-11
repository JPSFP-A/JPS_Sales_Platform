import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

prs = Presentation(r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx')
master = prs.slide_masters[0]

# Dump theme XML
theme_part = master.part.part_related_by(
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
theme_xml = etree.tostring(etree.fromstring(theme_part.blob), pretty_print=True).decode()
print('=== THEME XML (colors + fonts only) ===')
import re
# Find color scheme
m = re.search(r'<a:clrScheme.*?</a:clrScheme>', theme_xml, re.DOTALL)
if m: print(m.group())
m = re.search(r'<a:fontScheme.*?</a:fontScheme>', theme_xml, re.DOTALL)
if m: print(m.group())

# Master background
print('\n=== MASTER BACKGROUND ===')
bg_xml = etree.tostring(master.background._element, pretty_print=True).decode()
print(bg_xml[:800])

# Check slide 2 (TOC) for hyperlinks
print('\n=== SLIDE 2 HYPERLINKS ===')
slide2 = prs.slides[1]
slide_xml = etree.tostring(slide2._element, pretty_print=True).decode()
hlins = re.findall(r'r:id="([^"]+)"[^>]*?type="[^"]*hyperlink[^"]*"', slide_xml)
print(f'Hyperlink rIds: {hlins}')
# Also check relationships
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
try:
    for rel in slide2.part.rels.values():
        if 'hyperlink' in rel.reltype.lower():
            print(f'  hyperlink rel: {rel.reltype} -> {rel._target}')
except Exception as e:
    print(f'  rels error: {e}')
# Check raw XML for hlinkClick
hlinks = re.findall(r'<a:hlinkClick[^/]*/>', slide_xml)
print(f'hlinkClick elements: {hlinks[:10]}')

# Slide number state across slides
print('\n=== SLIDE NUMBER FOOTER STATE ===')
for i, slide in enumerate(prs.slides[:5]):
    sxml = etree.tostring(slide._element, pretty_print=True).decode()
    hf = re.search(r'<p:hf[^>]*/>', sxml)
    print(f'  Slide {i+1} hf: {hf.group() if hf else "none"}')
