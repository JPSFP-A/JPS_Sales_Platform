import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'C:\Users\jwilson\Downloads\2026 Brand Refresh Template - Aptos Core Font.pptx')
print(f'Slide size: {prs.slide_width.inches:.3f}" x {prs.slide_height.inches:.3f}"')
print(f'Slide count: {len(prs.slides)}')
print(f'Layout count: {len(prs.slide_layouts)}')
print()

for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name if slide.slide_layout else 'N/A'
    texts = [s.text_frame.text[:80].strip() for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    shape_info = []
    for s in slide.shapes:
        stype = s.shape_type
        if stype == MSO_SHAPE_TYPE.PICTURE:
            shape_info.append(f'IMG:{s.name}')
        elif stype == 1:  # auto shape
            shape_info.append(f'RECT:{s.name}')
        elif s.has_text_frame:
            shape_info.append(f'TB:{s.name}')
    bg = slide.background
    fill = bg.fill
    print(f'Slide {i+1}: layout={layout_name!r}  bg_fill={fill.type}')
    print(f'  Texts: {texts[:4]}')
    print(f'  Shapes: {shape_info[:8]}')
    print()
