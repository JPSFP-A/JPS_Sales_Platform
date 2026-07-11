"""
Apply JPS 2026 Brand to May 2026 Corporate Performance Report.

Strategy:
  - Cover              → full rebuild (gradient bg, large date/title text)
  - Section dividers   → full rebuild (gradient bg, Cyber Yellow accent bar)
  - Thank You          → full rebuild (gradient + centered large logo + white strip)
  - Content slides     → white bg + Electric Blue accent bar + logo + footer text
                         (NO Continuous Line — matches official template)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SRC   = r'C:\Users\jwilson\Downloads\May 2026 Corporate Performance Report.pptx'
OUT   = r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx'
BRAND = r'D:\Projects\Sales_Platform\analysis\brand_ref_2026\Brand Reference 2026 - Copy'
MEDIA = r'D:\Projects\Sales_Platform\analysis\template_media'

# Brand colours
OX = RGBColor(0x06, 0x25, 0x52)   # Oxford Blue
RY = RGBColor(0x0E, 0x4E, 0x95)   # Royal Blue
EL = RGBColor(0x00, 0x9F, 0xDA)   # Electric Blue
YL = RGBColor(0xFF, 0xF0, 0x00)   # Cyber Yellow
WH = RGBColor(0xFF, 0xFF, 0xFF)   # White

# Slide geometry (13.333" × 7.5")
SW = Inches(13.333)
SH = Inches(7.5)

# Small logo top-right (content slides, sections)
LW, LH = Inches(1.55), Inches(0.527)
LX, LY = SW - Inches(0.2) - LW, Inches(0.1)

# Continuous Line brand mark (1920×93, ratio 20.645:1)
BMH = Inches(13.333 / 20.645)    # ≈ 0.646"
BMY = SH - BMH                    # bottom-flush

# Template-exact accent bar dimensions
# Content slides — Electric Blue, from layout Freeform 5
BAR_CX, BAR_CY = Inches(0.599), Inches(0.674)
BAR_CW, BAR_CH = Inches(0.175), Inches(0.872)

# Section dividers — Cyber Yellow, from layout Freeform 5
BAR_SX, BAR_SY = Inches(0.458), Inches(2.19)
BAR_SW, BAR_SH = Inches(0.257), Inches(3.12)

# Footer position (from template TextBox 8)
FTR_X, FTR_Y = Inches(0.513), Inches(7.022)
FTR_W, FTR_H = Inches(7.241), Inches(0.317)
FOOTER_TEXT   = 'Corporate Performance Report  ·  May 2026'

# Asset paths
LOGO     = os.path.join(BRAND, 'JPS-Logo-Primary-FullColour.png')
BM_WHITE = os.path.join(BRAND, 'JPS-BrandMark-(White).png')
BM_GRAD  = os.path.join(BRAND, 'JPS-BrandMark-(Oxford-Gradient).png')
BG       = os.path.join(MEDIA, 'bg_gradient.png')

# Slide catalogue (0-indexed)
COVER_IDX  = 0
THANKS_IDX = 33
SECTION = {
    3:  'Financial\nHighlights',
    10: 'Operational\nExpenditure',
    14: 'Capital\nExpenditure',
    16: 'Cash Flow and\nFunding Outlook',
    19: 'Appendices',
}

# ── Helpers ───────────────────────────────────────────────────────────────────

def z_back(slide, shape):
    tree = slide.shapes._spTree
    el   = shape._element
    tree.remove(el)
    tree.insert(2, el)

def clear_shapes(slide):
    tree = slide.shapes._spTree
    tags = {'sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'}
    for el in [e for e in list(tree) if e.tag.split('}')[-1] in tags]:
        tree.remove(el)

def set_white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WH

def pic(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, x, y, w, h)

def tb(slide, text, x, y, w, h,
       face='Aptos', size=24, bold=False, italic=False,
       color=WH, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = face
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return box

def rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb      = fill
    shp.line.width          = 0
    return shp

def add_logo_small(slide):
    return pic(slide, LOGO, LX, LY, LW, LH)

def add_bm(slide, dark=True):
    path = BM_WHITE if dark else BM_GRAD
    return pic(slide, path, Inches(0), BMY, SW, BMH)

# ── Slide builders ────────────────────────────────────────────────────────────

def build_cover(slide):
    clear_shapes(slide)
    bg = pic(slide, BG, 0, 0, SW, SH); z_back(slide, bg)
    add_logo_small(slide)

    tb(slide, 'FINANCE DIVISION  ·  CORPORATE PERFORMANCE REPORT',
       Inches(1), Inches(2.0), Inches(11.3), Inches(0.45),
       face='Aptos', size=12, bold=True, color=YL, align=PP_ALIGN.CENTER)

    tb(slide, 'MAY 2026',
       Inches(0.4), Inches(2.55), Inches(12.5), Inches(1.9),
       face='Aptos Black', size=84, bold=True, color=WH, align=PP_ALIGN.CENTER)

    tb(slide, 'CORPORATE PERFORMANCE SUMMARY',
       Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.7),
       face='Aptos Black', size=28, bold=True, color=WH, align=PP_ALIGN.CENTER)

    tb(slide, 'Vernon Douglas, Chief Financial Officer',
       Inches(1), Inches(5.1), Inches(11.3), Inches(0.38),
       face='Aptos', size=14, color=EL, align=PP_ALIGN.CENTER)

    tb(slide, 'Finance Division  ·  June 2026',
       Inches(1), Inches(5.48), Inches(11.3), Inches(0.32),
       face='Aptos', size=12, color=WH, align=PP_ALIGN.CENTER)

    # Powering What Matters + Continuous Line (cover only gets the brand mark)
    tb(slide, 'Powering What Matters',
       Inches(7.685), Inches(6.918), Inches(5.0), Inches(0.31),
       face='Aptos', size=11, bold=True, color=WH, align=PP_ALIGN.RIGHT)
    add_bm(slide, dark=True)


def build_section(slide, title):
    clear_shapes(slide)
    bg = pic(slide, BG, 0, 0, SW, SH); z_back(slide, bg)
    add_logo_small(slide)

    # Cyber Yellow accent bar — template exact: x=0.458" y=2.19" w=0.257" h=3.12"
    rect(slide, BAR_SX, BAR_SY, BAR_SW, BAR_SH, YL)

    tb(slide, title,
       Inches(0.917), Inches(2.19), Inches(11.5), Inches(3.12),
       face='Aptos Black', size=46, bold=True, color=WH, align=PP_ALIGN.LEFT)

    # Footer text (white on dark)
    tb(slide, FOOTER_TEXT,
       FTR_X, FTR_Y, FTR_W, FTR_H,
       face='Aptos', size=11, bold=True, color=WH, align=PP_ALIGN.LEFT)


def build_thanks(slide):
    """Thank You — gradient + large centred logo + white bottom strip with Oxford-Gradient brand mark."""
    clear_shapes(slide)

    # Gradient covers most of slide (top 6.65"), white strip below
    bg = pic(slide, BG, 0, 0, SW, Inches(6.65)); z_back(slide, bg)
    # White strip at bottom
    rect(slide, Inches(0), Inches(6.65), SW, Inches(0.85), WH)

    # JPS logo — large, centred (template: w=3.038" h≈1.034", y=0.694", centred)
    big_lw, big_lh = Inches(3.038), Inches(1.034)
    big_lx = (SW - big_lw) / 2
    pic(slide, LOGO, big_lx, Inches(0.694), big_lw, big_lh)

    # THANK YOU — template: x=3.808" y=2.51" w=5.767" h=2.661" size=97pt
    tb(slide, 'THANK YOU',
       Inches(3.808), Inches(2.51), Inches(5.767), Inches(2.661),
       face='Aptos Black', size=97, bold=True, color=WH, align=PP_ALIGN.CENTER)

    # Oxford-Gradient brand mark on white strip
    add_bm(slide, dark=False)

    # Powering What Matters in Oxford Blue (on the white strip)
    tb(slide, 'Powering What Matters',
       Inches(7.685), Inches(6.918), Inches(5.0), Inches(0.31),
       face='Aptos', size=11, bold=True, color=OX, align=PP_ALIGN.RIGHT)


def brand_content(slide):
    """
    Brand an existing content slide.
    - Force white background
    - Add Electric Blue accent bar (template exact position)
    - Add JPS logo top-right
    - Add footer text
    - NO Continuous Line brand mark (matches official template)
    """
    set_white_bg(slide)
    add_logo_small(slide)
    # Electric Blue accent bar — template exact: x=0.599" y=0.674" w=0.175" h=0.872"
    rect(slide, BAR_CX, BAR_CY, BAR_CW, BAR_CH, EL)
    # Footer
    tb(slide, FOOTER_TEXT,
       FTR_X, FTR_Y, FTR_W, FTR_H,
       face='Aptos', size=11, bold=True, color=OX, align=PP_ALIGN.LEFT)


# ── Main ──────────────────────────────────────────────────────────────────────
prs = Presentation(SRC)

for i, slide in enumerate(prs.slides):
    label = f'slide {i+1}'
    if i == COVER_IDX:
        build_cover(slide);              print(f'{label}: cover rebuilt')
    elif i == THANKS_IDX:
        build_thanks(slide);             print(f'{label}: thank-you rebuilt')
    elif i in SECTION:
        build_section(slide, SECTION[i]); print(f'{label}: section — {SECTION[i]!r}')
    else:
        brand_content(slide);            print(f'{label}: content branded')

prs.save(OUT)
print(f'\nSaved -> {OUT}')
print(f'Size: {round(os.path.getsize(OUT)/1024)}KB')
