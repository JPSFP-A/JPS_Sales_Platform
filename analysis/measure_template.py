import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'C:\Users\jwilson\Downloads\2026 Brand Refresh Template - Aptos Core Font.pptx')

def rgb_str(color):
    try:
        return f'#{color.rgb}'
    except:
        return 'theme/inherit'

for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f'\n=== Slide {i+1}: {layout_name} ===')
    for s in slide.shapes:
        x = round(s.left / 914400, 3) if s.left else 0
        y = round(s.top / 914400, 3) if s.top else 0
        w = round(s.width / 914400, 3) if s.width else 0
        h = round(s.height / 914400, 3) if s.height else 0
        stype = s.shape_type
        fill_info = ''
        if stype == 1:  # auto shape
            try:
                fill_info = f' fill={rgb_str(s.fill.fore_color)}'
            except:
                fill_info = f' fill=type({s.fill.type})'
        text_info = ''
        if s.has_text_frame:
            txt = s.text_frame.text[:50].strip()
            if txt:
                try:
                    run = s.text_frame.paragraphs[0].runs[0]
                    clr = rgb_str(run.font.color)
                    sz = run.font.size // 12700 if run.font.size else 'inherit'
                    bold = run.font.bold
                    text_info = f' txt={txt!r} color={clr} size={sz} bold={bold}'
                except:
                    text_info = f' txt={txt!r}'
        print(f'  {s.name}: x={x}" y={y}" w={w}" h={h}"{fill_info}{text_info}')

    # Also check layout shapes
    layout = slide.slide_layout
    print(f'  [LAYOUT shapes: {layout.name}]')
    for s in layout.shapes:
        x = round(s.left / 914400, 3) if s.left else 0
        y = round(s.top / 914400, 3) if s.top else 0
        w = round(s.width / 914400, 3) if s.width else 0
        h = round(s.height / 914400, 3) if s.height else 0
        fill_info = ''
        if s.shape_type == 1:
            try:
                fill_info = f' fill={rgb_str(s.fill.fore_color)}'
            except:
                fill_info = f' fill=type({s.fill.type})'
        print(f'    {s.name}: x={x}" y={y}" w={w}" h={h}"{fill_info}')
