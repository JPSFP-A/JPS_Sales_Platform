import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

prs = Presentation(r'D:\Projects\Sales_Platform\analysis\May_2026_Corp_Perf_JPS_Brand.pptx')
slide = prs.slides[1]
slide_part = slide.part

# Map slide rIds to slide numbers
print('=== Slide part relationships ===')
for rid, rel in slide_part.rels.items():
    print(f'  {rid}: {rel.reltype.split("/")[-1]} -> {rel._target}')

# Find hlinkClick rIds in TOC
print('\n=== TOC hlinkClick rIds ===')
for shape in slide.shapes:
    if 'Content' in shape.name and shape.has_text_frame:
        for i, para in enumerate(shape.text_frame.paragraphs):
            for run in para.runs:
                rpr = run._r.find(f'{{{A}}}rPr')
                if rpr is not None:
                    hc = rpr.find(f'{{{A}}}hlinkClick')
                    if hc is not None:
                        rid = hc.get(f'{{{R}}}id')
                        action = hc.get('action', '')
                        print(f'  Para {i} run "{run.text}": rid={rid} action={action}')

# Also build slide id map
print('\n=== Slide index -> rId in presentation ===')
prs_part = prs.part
for i, slide_s in enumerate(prs.slides):
    # Get the slide's part name
    print(f'  Slide {i+1}: {slide_s.part.partname}')
